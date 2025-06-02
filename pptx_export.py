import tempfile
from flask import send_file, request, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from content_utils import process_content_for_layout


import textwrap

def hex_to_rgb(hex_color):
    """Convert hex color string to RGB tuple for PowerPoint"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    

def parse_template_colors(template_id):
    """Parse template colors for PowerPoint"""
    templates = {
        'corporate': {
            'colors': {
                'primary': '#0f4c81',
                'secondary': '#6e9cc4',
                'accent': '#f2b138',
                'background': '#ffffff',
                'text': '#333333'
            },
            'font': 'Arial'
        },
        'creative': {
            'colors': {
                'primary': '#ff6b6b',
                'secondary': '#4ecdc4',
                'accent': '#ffd166',
                'background': '#f9f1e6',
                'text': '#5a3921'
            },
            'font': 'Georgia'
        },
        'creative updated': {
        'colors': {
            'primary': '#3A86FF',
            'secondary': '#8338EC',
            'accent': '#00CFC1',
            'background': '#FDFDFD',
            'text': '#1D1D1D'
        },
        'font': 'Inter'
    },
        'minimal': {
            'colors': {
                'primary': '#2c3e50',
                'secondary': '#95a5a6',
                'accent': '#e74c3c',
                'background': '#f8f8f8',
                'text': '#222222'
            },
            'font': 'Helvetica'
        },
        'nature': {
            'colors': {
                'primary': '#2d6a4f',
                'secondary': '#74c69d',
                'accent': '#ffb703',
                'background': '#f1faee',
                'text': '#1b4332'
            },
            'font': 'Verdana'
        },
        'tech': {
            'colors': {
                'primary': '#3a0ca3',
                'secondary': '#4361ee',
                'accent': '#7209b7',
                'background': '#f8f9fa',
                'text': '#212529'
            },
            'font': 'Roboto'
        },
        'gradient': {
            'colors': {
                'primary': '#8338ec',
                'secondary': '#3a86ff',
                'accent': '#ff006e',
                'background': '#ffffff',
                'text': '#2b2d42'
            },
            'font': 'Segoe UI'
        },
        'pastel': {
            'colors': {
                'primary': '#9d8189',
                'secondary': '#d8e2dc',
                'accent': '#ffcad4',
                'background': '#f7ede2',
                'text': '#6d6875'
            },
            'font': 'Trebuchet MS'
        },
        'energetic': {
            'colors': {
                'primary': '#f72585',
                'secondary': '#7209b7',
                'accent': '#4cc9f0',
                'background': '#ffffff',
                'text': '#2b2d42'
            },
            'font': 'Tahoma'
        },
        'earthy': {
            'colors': {
                'primary': '#996633',
                'secondary': '#dda15e',
                'accent': '#bc6c25',
                'background': '#fefae0',
                'text': '#283618'
            },
            'font': 'Palatino'
        },
        'ocean': {
            'colors': {
                'primary': '#003566',
                'secondary': '#468faf',
                'accent': '#00b4d8',
                'background': '#f0f7f9',
                'text': '#001845'
            },
            'font': 'Calibri'
        },
        'monochrome': {
            'colors': {
                'primary': '#2b2d42',
                'secondary': '#8d99ae',
                'accent': '#ef233c',
                'background': '#edf2f4',
                'text': '#1a1a1a'
            },
            'font': 'Courier New'
        },
        'sunset': {
            'colors': {
                'primary': '#ff9e00',
                'secondary': '#ff4d00',
                'accent': '#7678ed',
                'background': '#fff1e6',
                'text': '#3d405b'
            },
            'font': 'Garamond'
        },
        'botanical': {
            'colors': {
                'primary': '#386641',
                'secondary': '#a7c957',
                'accent': '#fb8500',
                'background': '#f8f9fa',
                'text': '#283618'
            },
            'font': 'Lucida Sans'
        }
    }

    
    template = templates.get(template_id, templates['corporate'])
    
    # Convert hex colors to RGB
    rgb_colors = {}
    for key, hex_value in template['colors'].items():
        rgb_colors[key] = hex_to_rgb(hex_value)
    
    return {
        'colors': rgb_colors,
        'font': template['font']
    }

def apply_template_to_slide(slide, template):
    """Apply template styling to a slide"""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = template['colors']['background']
    return template

def apply_text_formatting(paragraph, font_size, color, font_name, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    """Apply consistent text formatting to a paragraph"""
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.name = font_name
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.alignment = alignment

def clean_text_content(content):
    """Clean and normalize text content from JSON for PPTX"""
    if isinstance(content, dict) and 'text' in content:
        return content['text']
    elif isinstance(content, list):
        return "\n".join([str(item) for item in content])
    return str(content or '')

def create_image_with_features_slide(presentation, content, template):
    """Create a slide with an image and feature points with circular icons that match the preview"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Add title
    left = Inches(0.1)
    top = Inches(0)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    
    title_text = content.get('title', 'Features')
    if isinstance(title_text, dict) and 'text' in title_text:
        title_text = title_text['text']
    title_text = str(title_text)
    
    p = text_frame.add_paragraph()
    p.text = title_text
    apply_text_formatting(
        p, 
        font_size=22,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Add image placeholder on the left
    left = Inches(0.2)
    top = Inches(1)
    width = Inches(4)
    height = Inches(4.5)
    
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # Using proper enumeration for rectangle
        left, top, width, height
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = template['colors']['secondary']
    img_placeholder.line.color.rgb = template['colors']['primary']
    
    # FIXED APPROACH: Add image description text with explicit positioning 
    # that doesn't depend on other variables
    img_text_left = Inches(0.2)  # Fixed position in the center of image area
    img_text_top = Inches(0.8)  # Fixed position in the center of image area
    img_text_width = Inches(4.0)  # Fixed width
    img_text_height = Inches(0.6)  # Fixed height
    
    
    # Add the text on top of the background
    img_text_box = slide.shapes.add_textbox(
        img_text_left, img_text_top, img_text_width, img_text_height
    )
    text_frame = img_text_box.text_frame
    text_frame.word_wrap = True
    
    img_desc = content.get('imageDescription', 'Image')
    if isinstance(img_desc, dict) and 'text' in img_desc:
        img_desc = img_desc['text']
    img_desc = str(img_desc)
    
    if "**Image Prompt:**" in img_desc:
        img_desc = img_desc.split("**Image Prompt:**", 1)[1].strip()
    img_desc = img_desc.replace('*', '')

    p = text_frame.add_paragraph()
    p.text = img_desc
    apply_text_formatting(
        p, 
        font_size=8,
        color=RGBColor(255, 255, 255),
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT  # Changed to LEFT alignment
    )
        
    # Add features on the right side with proper circular icons
    features = content.get('features', [])
    
    for i, feature in enumerate(features[:4]):  # Limit to 4 features
        feature_top = Inches(1.5 + i * 0.9)
        
        # Add icon/circle with proper positioning
        icon_left = Inches(4.8)
        icon_top = feature_top - Inches(0.3)
        icon_size = Inches(0.5)
        
        # KEY FIX: Use MSO_SHAPE enumeration instead of hardcoded number
        # and ensure width = height for a perfect circle
        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,  # Proper enumeration for OVAL shape
            icon_left, icon_top, icon_size, icon_size  # Equal width and height
        )
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = template['colors']['primary']
        
        # ADD GEAR ICON ON THE CIRCLE
        gear_text_box = slide.shapes.add_textbox(
            icon_left, 
            icon_top - Inches(0.15),
            icon_size, 
            icon_size
        )
        gear_text_frame = gear_text_box.text_frame
        gear_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        gear_p = gear_text_frame.add_paragraph()
        gear_p.text = "⚙️"  # Gear emoji
        gear_p.alignment = PP_ALIGN.CENTER
        
        # Apply formatting to gear icon text
        for run in gear_p.runs:
            run.font.size = Pt(16)  # Adjust size as needed
            run.font.color.rgb = RGBColor(255, 255, 255)  # White color
        
        # Add feature title - ALIGNED EXACTLY WITH ICON VERTICAL POSITION
        title_left = Inches(5.3)  # Moved closer to the icon
        title_width = Inches(4.0)  # Slightly wider to accommodate more text
        title_height = Inches(0.2)
        
        # The key change: Use the exact same vertical position as the icon
        # Calculate vertical center alignment with the icon
        title_top = icon_top + (icon_size/2) - (title_height/2) - Inches(0.3)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        text_frame = title_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center the text in the box
        
        feature_title = feature.get('title', f'Feature {i+1}')
        if isinstance(feature_title, dict) and 'text' in feature_title:
            feature_title = feature_title['text']
        feature_title = str(feature_title)
        
        p = text_frame.add_paragraph()
        p.text = feature_title
        apply_text_formatting(
            p, 
            font_size=16,
            color=template['colors']['primary'],
            font_name=template['font'],
            bold=True
        )
        
        # Add feature description - POSITIONED RELATIVE TO TITLE BOX
        desc_top = title_top + title_height - Inches(0.1) # Position directly below the title box
        desc_box = slide.shapes.add_textbox(title_left, desc_top, title_width, Inches(0.5))
        text_frame = desc_box.text_frame
        text_frame.word_wrap = True
        
        feature_desc = feature.get('description', '')
        if isinstance(feature_desc, dict) and 'text' in feature_desc:
            feature_desc = feature_desc['text']
        feature_desc = str(feature_desc)
        
        p = text_frame.add_paragraph()
        p.text = feature_desc
        apply_text_formatting(
            p, 
            font_size=10,
            color=template['colors']['text'],
            font_name=template['font']
        )
    return slide

def create_presentation(filename, slides, template_id):
    """Create a PowerPoint presentation with improved layout matching"""
    prs = Presentation()
    
    # Set the slide size to 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Parse template colors
    template = parse_template_colors(template_id)
    
    # Create slides based on their layout type
    for slide_data in slides:
        layout = slide_data.get('layout')
        content = slide_data.get('content', {})
        
        if layout == 'titleOnly':
            create_title_only_slide(prs, content, template)
        elif layout == 'titleAndBullets':
            create_title_and_bullets_slide(prs, content, template)
        elif layout == 'quote':
            create_quote_slide(prs, content, template)
        elif layout == 'imageAndParagraph':
            create_image_and_paragraph_slide(prs, content, template)
        elif layout == 'twoColumn':
            create_two_column_slide(prs, content, template)
        elif layout == 'imageWithFeatures':
            create_image_with_features_slide(prs, content, template)
        elif layout == 'numberedFeatures':
            create_numbered_features_slide(prs, content, template)
        elif layout == 'benefitsGrid':
            create_benefits_grid_slide(prs, content, template)
        elif layout == 'iconGrid':
            create_icon_grid_slide(prs, content, template)
        elif layout == 'sideBySideComparison':
            create_side_by_side_comparison_slide(prs, content, template)
        elif layout == 'timeline':
            create_timeline_slide(prs, content, template)
        elif layout == 'conclusion':
            create_conclusion_slide(prs, content, template)
    
    # Save the presentation
    prs.save(filename)
    
    return filename


def create_title_only_slide(presentation, content, template):
    """Create a title-only slide that matches the preview slide exactly"""
    slide_layout = presentation.slide_layouts[0]  # Title Slide layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Get title text
    title_text = clean_text_content(content.get('title', 'Title'))
    subtitle_text = clean_text_content(content.get('subtitle', 'Subtitle'))
    
    # Create a title textbox with adjusted positioning (reduced space from top)
    left = Inches(1.0)
    top = Inches(0.5)  # Reduced from 2.0 to 1.2 inches from top
    width = Inches(9.0)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    
    # Add title text with smaller font size
    p = text_frame.add_paragraph()
    p.text = title_text
    apply_text_formatting(
        p, 
        font_size=34,  # Reduced from 54 to 44
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Add subtitle with adjusted positioning
    subtitle_left = Inches(1.0)
    subtitle_top = Inches(2.5)  # Reduced from 3.8 to 3.0
    subtitle_width = Inches(8.0)
    subtitle_height = Inches(1.0)
    
    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    text_frame = subtitle_box.text_frame
    text_frame.word_wrap = True
    
    # Add subtitle text with smaller font
    p = text_frame.add_paragraph()
    p.text = subtitle_text
    apply_text_formatting(
        p, 
        font_size=24,  # Reduced from 32 to 26
        color=template['colors']['secondary'],
        font_name=template['font'],
        alignment=PP_ALIGN.CENTER
    )
    
    return slide



#updated responsive fucntions 
def create_timeline_slide(presentation, content, template):
    """Create a slide with a chronological timeline with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    # Apply template
    template = apply_template_to_slide(slide, template)

    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Timeline'))
    events = content.get('events', [])

    # Calculate dynamic font size for title
    title_length = len(title)
    title_font_size = max(20, min(28, int(28 * (40 / max(40, title_length)))))

    # Add title - REDUCED SIZE AND MOVED UP
    left = Inches(0.1)
    top = Inches(-0.2)  # Moved up above slide top edge
    width = Inches(9)
    height = Inches(0.6)  # Reduced height

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add timeline events with vertical connector
    # Timeline starts at higher vertical position
    timeline_start_y = Inches(0.7)  # Reduced from 1.5
    timeline_x = Inches(1.6)  # X position of vertical line

    # Add vertical timeline line first as background - REDUCED HEIGHT
    if len(events) > 1:
        # Reduce spacing between events from 1.0 to 0.8 inches
        line_height = Inches(0.7 + (len(events) - 1) * 0.8)  # Further reduced spacing
        timeline_line = slide.shapes.add_shape(
            1,  # Rectangle
            timeline_x - Inches(0.01),  # Center line
            timeline_start_y + Inches(0.2),  # Start below first circle
            Inches(0.02),  # Very thin line
            line_height
        )
        timeline_line.fill.solid()
        timeline_line.fill.fore_color.rgb = template['colors']['secondary']

    for i, event in enumerate(events):
        event_y = timeline_start_y + Inches(i * 0.8)
        
        # Get event content with no truncation
        year = clean_text_content(event.get('year', ''))
        event_title = clean_text_content(event.get('title', f'Event {i+1}'))
        description = clean_text_content(event.get('description', ''))
        
        # Calculate dynamic font sizes based on content length
        year_length = len(year)
        event_title_length = len(event_title)
        description_length = len(description)
        
        year_font_size = max(10, min(14, int(14 * (10 / max(10, year_length)))))
        event_title_font_size = max(12, min(16, int(16 * (30 / max(30, event_title_length)))))
        description_font_size = max(10, min(12, int(12 * (100 / max(100, description_length)))))
        
        # HARDCODED CIRCLE - using MSO_SHAPE.OVAL for a perfect circle
        circle_size = Inches(0.1)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            timeline_x - circle_size/2,  # Center horizontally on timeline
            event_y + Inches(0.1),  # Align vertically with text
            circle_size,
            circle_size  # Equal width and height for perfect circle
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = template['colors']['accent']
        
        # Add year label attached to circle but moved upside
        year_label = slide.shapes.add_textbox(
            timeline_x - Inches(1.1),  # Position left of circle
            event_y - Inches(0.2),  # Moved up above the circle
            Inches(1.0),
            Inches(0.15)  
        )
        text_frame = year_label.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = text_frame.add_paragraph()
        p.text = year
        apply_text_formatting(
            p, 
            font_size=year_font_size,
            color=template['colors']['secondary'],
            font_name=template['font'],
            bold=True,
            alignment=PP_ALIGN.RIGHT  # Right align text to align with timeline
        )
        
        # Add event title to the right of timeline
        title_left = timeline_x + Inches(0.2)  # Space after timeline
        title_box = slide.shapes.add_textbox(
            title_left,
            event_y - Inches(0.3),
            Inches(7.5),  # Width for text
            Inches(0.3)  # Reduced from 0.4
        )
        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = event_title
        apply_text_formatting(
            p, 
            font_size=event_title_font_size,
            color=template['colors']['primary'],
            font_name=template['font'],
            bold=True
        )
        
        # Add event description below title
        desc_box = slide.shapes.add_textbox(
            title_left,
            event_y,  # Positioned directly below title
            Inches(8),  # Width for text
            Inches(0.6)  # Increased height for longer content
        )
        text_frame = desc_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Add a period at the end of the description if it doesn't already end with one
        if description and not description.endswith('.'):
            description += '.'
        
        p = text_frame.add_paragraph()
        p.text = description
        apply_text_formatting(
            p, 
            font_size=description_font_size,
            color=template['colors']['text'],
            font_name=template['font']
        )

    return slide

def create_conclusion_slide(presentation, content, template):
    """Create a conclusion slide with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    # Apply template
    template = apply_template_to_slide(slide, template)

    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Key Takeaways'))
    summary = clean_text_content(content.get('summary', ''))
    nextSteps = content.get('nextSteps', [])

    # Calculate dynamic font sizes based on content length
    title_length = len(title)
    summary_length = len(summary)

    title_font_size = max(16, min(20, int(20 * (40 / max(40, title_length)))))
    summary_font_size = max(12, min(16, int(16 * (150 / max(150, summary_length)))))

    # Add title
    left = Inches(0.1)
    top = Inches(0.1)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    # Enable text wrapping for title
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add summary paragraph
    summary_top = Inches(1)
    summary_box = slide.shapes.add_textbox(left, summary_top, width, Inches(1.5))
    text_frame = summary_box.text_frame
    # Enable text wrapping for summary
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = summary
    apply_text_formatting(
        p, 
        font_size=summary_font_size,
        color=template['colors']['secondary'],
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT
    )

    # Add next steps if present
    if nextSteps:
        # Moved further down from 2 to 2.5 inches
        steps_top = Inches(2.1)
        steps_title_box = slide.shapes.add_textbox(left, steps_top, width, Inches(0.5))
        text_frame = steps_title_box.text_frame
        
        p = text_frame.add_paragraph()
        p.text = "Next Steps"
        apply_text_formatting(
            p, 
            font_size=16,
            color=template['colors']['primary'],
            font_name=template['font'],
            bold=True,
            alignment=PP_ALIGN.LEFT
        )
        
        # Add the steps as bullet points
        steps_box = slide.shapes.add_textbox(left, steps_top + Inches(0.7), width, Inches(1.8))
        text_frame = steps_box.text_frame
        # Enable text wrapping for bullet points
        text_frame.word_wrap = True
        
        for i, step in enumerate(nextSteps):
            step_text = clean_text_content(step)
            
            # Calculate dynamic font size for step
            step_length = len(step_text)
            step_font_size = max(10, min(12, int(12 * (50 / max(50, step_length)))))
            
            p = text_frame.add_paragraph()
            p.text = '• ' + step_text
            apply_text_formatting(
                p, 
                font_size=step_font_size,
                color=template['colors']['text'],
                font_name=template['font'],
                alignment=PP_ALIGN.LEFT
            )
            p.level = 0
            
            # Add space between pointers by adding space after each paragraph
            p.space_after = Pt(12)  # Adding 12 points of space after each bullet point

    return slide

def create_numbered_features_slide(presentation, content, template):
    """Create a slide with numbered features with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    # Apply template
    template = apply_template_to_slide(slide, template)

    # Get content with no truncation - ADD DEFAULT VALUES
    title = clean_text_content(content.get('title', 'Key Features'))  # Fixed with default
    features = content.get('features', [])  # This is fine

    # Calculate dynamic font size for title
    title_length = len(title)
    title_font_size = max(20, min(28, int(28 * (60 / max(60, title_length)))))

    # Add title
    left = Inches(0.2)
    top = Inches(0.1)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add features in a grid with adaptive sizing
    grid_width = 2

    for i, feature in enumerate(features[:4]):  # Limit to 4 features
        row = i // grid_width
        col = i % grid_width
        
        feature_left = Inches(0.5 + col * 4.5)
        feature_top = Inches(1.3 + row * 1.8)
        
        # Get feature content with no truncation
        feature_number = feature.get('number', str(i + 1))
        feature_title = clean_text_content(feature.get('title', f'Feature {i+1}'))
        feature_description = clean_text_content(feature.get('description', f'Description {i+1}'))
        
        # Calculate dynamic font sizes for feature content
        feature_title_length = len(feature_title)
        feature_desc_length = len(feature_description)
        
        feature_title_font_size = max(14, min(16, int(16 * (40 / max(40, feature_title_length)))))
        feature_desc_font_size = max(10, min(14, int(14 * (100 / max(100, feature_desc_length)))))
        
        # Add number circle
        number_size = Inches(0.6)
        number_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            feature_left, feature_top, number_size, number_size
        )
        number_shape.fill.solid()
        number_shape.fill.fore_color.rgb = template['colors']['secondary']
        
        # Add number text
        number_text_box = slide.shapes.add_textbox(
            feature_left,
            feature_top - Inches(0.15),
            number_size,
            number_size
        )
        text_frame = number_text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = str(feature_number)
        apply_text_formatting(
            p, 
            font_size=18,
            color=template['colors']['background'],
            font_name=template['font'],
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        
        # Add feature title with adaptive sizing
        title_left = feature_left + Inches(0.8)
        title_width = Inches(3.5)
        title_height = Inches(0.4)
        
        # Position it slightly higher
        title_top_offset = Inches(-0.3)
        
        title_box = slide.shapes.add_textbox(
            title_left, 
            feature_top + title_top_offset,
            title_width, 
            title_height
        )        
        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = feature_title
        apply_text_formatting(
            p, 
            font_size=feature_title_font_size,
            color=template['colors']['primary'],
            font_name=template['font'],
            bold=True
        )
        
        # Add feature description with adaptive sizing
        desc_top_offset = Inches(0.28)
        desc_box = slide.shapes.add_textbox(
            title_left, 
            feature_top + title_top_offset + desc_top_offset,
            title_width, 
            Inches(0.8)  # Increased height for longer content
        )
        text_frame = desc_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = feature_description
        apply_text_formatting(
            p, 
            font_size=feature_desc_font_size,
            color=template['colors']['text'],
            font_name=template['font']
        )

    return slide
def create_benefits_grid_slide(presentation, content, template):
    """Create a slide with a benefits grid with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    # Apply template
    template = apply_template_to_slide(slide, template)

    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Benefits'))
    imageDescription = clean_text_content(content.get('imageDescription', 'Image'))
    benefits = content.get('benefits', [])

    # Calculate dynamic font size for title
    title_length = len(title)
    title_font_size = max(20, min(28, int(28 * (60 / max(60, title_length)))))

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add image on the left
    img_left = Inches(0.5)
    img_top = Inches(1.3)
    img_width = Inches(3.8)
    img_height = Inches(4.0)

    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        img_left, img_top, img_width, img_height
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = template['colors']['secondary']
    img_placeholder.line.color.rgb = template['colors']['secondary']

    # Add image placeholder text with adaptive sizing
    img_text_left = img_left + Inches(0.2)
    img_text_top = img_top + Inches(0.2)
    img_text_width = img_width - Inches(0.4)
    img_text_height = img_height - Inches(0.4)

    # Calculate font size for image description
    img_desc_length = len(imageDescription)
    img_desc_font_size = max(10, min(14, int(14 * (400 / max(400, img_desc_length)))))

    img_text_box = slide.shapes.add_textbox(
        img_text_left, img_text_top, img_text_width, img_text_height
    )
    text_frame = img_text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.add_paragraph()
    p.text = imageDescription
    p.alignment = PP_ALIGN.CENTER

    for run in p.runs:
        run.font.size = Pt(img_desc_font_size)
        run.font.name = template['font']
        run.font.color.rgb = RGBColor(255, 255, 255)

    # Add benefits in a grid with adaptive sizing
    benefits_grid_left = Inches(4.8)
    benefits_grid_top = img_top
    benefits_grid_width = Inches(4.7)
    benefits_grid_height = img_height

    # Distribute benefits in a 2x2 grid
    for i, benefit in enumerate(benefits[:4]):
        row = i // 2
        col = i % 2
        
        # Calculate cell width and height
        cell_width = benefits_grid_width / 2 - Inches(0.15)  # Subtract gap
        cell_height = benefits_grid_height / 2 - Inches(0.15)  # Subtract gap
        
        # Calculate cell position
        cell_left = benefits_grid_left + col * (cell_width + Inches(0.3))
        cell_top = benefits_grid_top + row * (cell_height + Inches(0.3))
        
        # Get benefit content with no truncation
        benefit_title = clean_text_content(benefit.get('title', f'Benefit {i+1}'))
        benefit_description = clean_text_content(benefit.get('description', ''))
        
        # Calculate dynamic font sizes for benefit content
        benefit_title_length = len(benefit_title)
        benefit_desc_length = len(benefit_description)
        
        benefit_title_font_size = max(10, min(16, int(16 * (30 / max(30, benefit_title_length)))))
        benefit_desc_font_size = max(8, min(14, int(14 * (60 / max(60, benefit_desc_length)))))
        
        # Create benefit card
        benefit_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cell_left, cell_top, cell_width, cell_height
        )
        benefit_card.fill.solid()
        benefit_card.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        benefit_card.line.fill.solid()
        benefit_card.line.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray border
        
        # Add left border accent strip
        border_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            cell_left + Inches(0.15), cell_top + Inches(0.15), Inches(0.04), Inches(0.35)
        )
        border_accent.fill.solid()
        border_accent.fill.fore_color.rgb = template['colors']['secondary']
        border_accent.line.fill.solid()
        border_accent.line.fill.fore_color.rgb = template['colors']['secondary']
        
        # Add benefit title
        title_left = cell_left + Inches(0.25)
        title_top = cell_top + Inches(0.15)
        title_width = cell_width - Inches(0.3)
        title_height = Inches(0.3)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = benefit_title
        apply_text_formatting(
            p, 
            font_size=benefit_title_font_size,
            color=template['colors']['primary'],
            font_name=template['font'],
            bold=True
        )
        
        # Add benefit description
        desc_left = cell_left + Inches(0.2)
        desc_top = cell_top + Inches(0.5)
        desc_width = cell_width - Inches(0.3)
        desc_height = cell_height - Inches(0.6)
        
        desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
        text_frame = desc_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = benefit_description
        apply_text_formatting(
            p, 
            font_size=benefit_desc_font_size,
            color=template['colors']['text'],
            font_name=template['font']
        )

    return slide
def create_icon_grid_slide(presentation, content, template):
    """Create a slide with a grid of category icons with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    # Apply template
    template = apply_template_to_slide(slide, template)

    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Categories'))
    categories = content.get('categories', [])

    # Calculate dynamic font size for title
    title_length = len(title)
    title_font_size = max(20, min(26, int(26 * (60 / max(60, title_length)))))

    # Add title
    left = Inches(0.2)
    top = Inches(0)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=RGBColor(255, 111, 97),  # Use the coral/red color
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Calculate grid layout
    row_count = 2
    col_count = 4

    # Smaller cards
    card_width = Inches(2.2)
    card_height = Inches(2.0)

    # Adjust spacing to ensure even distribution
    horizontal_gap = Inches(0.2)
    vertical_gap = Inches(0.3)

    # Move everything up
    start_x = Inches(0.3)
    start_y = Inches(1)

    # Create cards with perfectly round circle icons
    for i, category in enumerate(categories[:8]):
        row = i // col_count
        col = i % col_count
        
        # Calculate card position with proper spacing
        card_left = start_x + col * (card_width + horizontal_gap)
        card_top = start_y + row * (card_height + vertical_gap)
        
        # Get category content with no truncation
        category_name = clean_text_content(category.get('name', f'Category {i+1}'))
        category_description = clean_text_content(category.get('description', ''))
        
        # Calculate dynamic font sizes for category content
        category_name_length = len(category_name)
        category_desc_length = len(category_description)
        
        category_name_font_size = max(10, min(16, int(16 * (25 / max(25, category_name_length)))))
        category_desc_font_size = max(8, min(12, int(12 * (40 / max(40, category_desc_length)))))
        
        # Add white card background
        card = slide.shapes.add_shape(
            1,  # Rectangle
            card_left, card_top, card_width, card_height
        )
        
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray border
        card.shadow.inherit = False
        card.shadow.visible = True
        card.shadow.blur_radius = Pt(5)
        card.shadow.distance = Pt(2)
        card.shadow.angle = 45
        card.shadow.color = RGBColor(200, 200, 200)
        
        # Add perfectly round circle icon
        icon_size = Inches(0.5)  # Perfect circle needs same width and height
        icon_left = card_left + (card_width - icon_size) / 2  # Center in card
        icon_top = card_top + Inches(0.2)  # Position from top
        
        # Use OVAL shape type with equal width and height for a perfect circle
        circle_icon = slide.shapes.add_shape(
            5,  # Oval/circle shape (5 is OVAL in pptx)
            icon_left, icon_top, icon_size, icon_size
        )
        circle_icon.fill.solid()
        circle_icon.fill.fore_color.rgb = RGBColor(130, 120, 240)  # Purple
        # Remove any border to ensure perfect circle
        circle_icon.line.fill.solid()
        circle_icon.line.fill.fore_color.rgb = RGBColor(130, 120, 240)  # Match fill color
        
        # Add a star inside circle
        star_size = Inches(0.25)
        star_left = icon_left + (icon_size - star_size) / 2
        star_top = icon_top + (icon_size - star_size) / 2

        star_box = slide.shapes.add_textbox(
            star_left, star_top, star_size, star_size
        )
        text_frame = star_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = False
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = text_frame.paragraphs[0]
        p.text = "★"  # Star character
        apply_text_formatting(
            p, 
            font_size=14,
            color=RGBColor(255, 215, 0),  # Gold color for star
            font_name="Arial",
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        
        # Add category name
        name_left = card_left
        name_top = icon_top + icon_size + Inches(0)
        name_width = card_width
        name_height = Inches(0.3)
        
        name_box = slide.shapes.add_textbox(name_left, name_top, name_width, name_height)
        text_frame = name_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = category_name
        apply_text_formatting(
            p, 
            font_size=category_name_font_size,
            color=RGBColor(255, 111, 97),  # Coral/red color
            font_name=template['font'],
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        
        # Add category description
        desc_left = card_left + Inches(0.1)
        desc_top = name_top + Inches(0.3)
        desc_width = card_width - Inches(0.2)
        desc_height = Inches(0.8)
        
        desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
        text_frame = desc_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = category_description
        apply_text_formatting(
            p, 
            font_size=category_desc_font_size,
            color=RGBColor(90, 57, 33),  # Brown/dark color
            font_name=template['font'],
            alignment=PP_ALIGN.CENTER
        )

    return slide
def create_side_by_side_comparison_slide(presentation, content, template):
    """Create a slide with side-by-side comparison with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    # Apply template
    template = apply_template_to_slide(slide, template)

    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Comparison'))
    leftTitle = clean_text_content(content.get('leftTitle', 'The Challenge'))
    rightTitle = clean_text_content(content.get('rightTitle', 'The Solution'))
    leftPoints = content.get('leftPoints', [])
    rightPoints = content.get('rightPoints', [])

    # Calculate dynamic font sizes
    title_length = len(title)
    left_title_length = len(leftTitle)
    right_title_length = len(rightTitle)

    title_font_size = max(18, min(24, int(24 * (60 / max(60, title_length)))))
    left_title_font_size = max(16, min(22, int(22 * (30 / max(30, left_title_length)))))
    right_title_font_size = max(16, min(22, int(22 * (30 / max(30, right_title_length)))))

    # Add title
    left = Inches(0.1)
    top = Inches(0)
    width = Inches(9)
    height = Inches(0.6)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add left side image card with transparent background
    left_img_left = Inches(0.2)
    left_img_top = Inches(1.2)
    left_img_width = Inches(4.3)
    left_img_height = Inches(1.5)

    left_img_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left_img_left, left_img_top, left_img_width, left_img_height
    )
    left_img_card.fill.solid()
    left_img_card.fill.fore_color.rgb = template['colors']['secondary']
    left_img_card.fill.transparency = 0.5  # 50% transparent
    left_img_card.line.color.rgb = template['colors']['primary']
    left_img_card.line.width = Pt(0)

    # Add "Challenge Image" text in the center of the left image card
    left_img_text_box = slide.shapes.add_textbox(
        left_img_left, left_img_top + Inches(0.35), left_img_width, Inches(0.4)
    )
    text_frame = left_img_text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.add_paragraph()
    p.text = "Challenge Image"
    p.alignment = PP_ALIGN.CENTER
    apply_text_formatting(
        p, 
        font_size=18,
        color=template['colors']['background'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Add right side image card
    right_img_left = Inches(5.2)
    right_img_top = Inches(1.2)
    right_img_width = Inches(4.3)
    right_img_height = Inches(1.5)

    right_img_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_img_left, right_img_top, right_img_width, right_img_height
    )
    right_img_card.fill.solid()
    right_img_card.fill.fore_color.rgb = template['colors']['accent']
    right_img_card.fill.transparency = 0  # 0% transparent
    right_img_card.line.width = Pt(0)

    # Add "Solution Image" text in the center of the right image card
    right_img_text_box = slide.shapes.add_textbox(
        right_img_left, right_img_top + Inches(0.35), right_img_width, Inches(0.4)
    )
    text_frame = right_img_text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.add_paragraph()
    p.text = "Solution Image"
    p.alignment = PP_ALIGN.CENTER
    apply_text_formatting(
        p, 
        font_size=18,
        color=template['colors']['background'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Add left side title
    left_title_left = Inches(0.2)
    left_title_top = Inches(2.5)
    left_title_width = Inches(4.3)
    left_title_height = Inches(0.4)

    left_title_box = slide.shapes.add_textbox(left_title_left, left_title_top, left_title_width, left_title_height)
    text_frame = left_title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = leftTitle
    apply_text_formatting(
        p, 
        font_size=left_title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add right side title
    right_title_left = Inches(5.2)
    right_title_top = Inches(2.5)
    right_title_width = Inches(4.3)
    right_title_height = Inches(0.4)

    right_title_box = slide.shapes.add_textbox(right_title_left, right_title_top, right_title_width, right_title_height)
    text_frame = right_title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = text_frame.add_paragraph()
    p.text = rightTitle
    apply_text_formatting(
        p, 
        font_size=right_title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Add left side bullet points
    left_points_left = Inches(0.2)
    left_points_top = Inches(2.8)
    left_points_width = Inches(4.3)
    left_points_height = Inches(2.0)

    left_points_box = slide.shapes.add_textbox(left_points_left, left_points_top, left_points_width, left_points_height)
    text_frame = left_points_box.text_frame
    text_frame.word_wrap = True

    for point in leftPoints:
        point_text = clean_text_content(point)
        
        # Calculate dynamic font size for point
        point_length = len(point_text)
        point_font_size = max(10, min(12, int(12 * (80 / max(80, point_length)))))
        
        p = text_frame.add_paragraph()
        p.text = "• " + point_text
        apply_text_formatting(
            p, 
            font_size=point_font_size,
            color=template['colors']['text'],
            font_name=template['font'],
            alignment=PP_ALIGN.LEFT
        )
        p.space_after = Pt(12)  # Add space after each bullet point

    # Add right side bullet points
    right_points_left = Inches(5.2)
    right_points_top = Inches(2.8)
    right_points_width = Inches(4.3)
    right_points_height = Inches(2.0)

    right_points_box = slide.shapes.add_textbox(right_points_left, right_points_top, right_points_width, right_points_height)
    text_frame = right_points_box.text_frame
    text_frame.word_wrap = True

    for point in rightPoints:
        point_text = clean_text_content(point)
        
        # Calculate dynamic font size for point
        point_length = len(point_text)
        point_font_size = max(10, min(12, int(12 * (80 / max(80, point_length)))))
        
        p = text_frame.add_paragraph()
        p.text = "• " + point_text
        apply_text_formatting(
            p, 
            font_size=point_font_size,
            color=template['colors']['text'],
            font_name=template['font'],
            alignment=PP_ALIGN.LEFT
        )
        p.space_after = Pt(12)  # Add space after each bullet point

    return slide
    ## 3. Updates to `pptx_export.py`
from pptx.enum.text import MSO_AUTO_SIZE

def create_title_and_bullets_slide(presentation, content, template):
    """Create a slide with title and bullet points that adapts to content length"""
    slide_layout = presentation.slide_layouts[6]  # Use blank layout for more control
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Title'))
    bullets = content.get('bullets', [])
    
    # Calculate title font size based on length
    title_length = len(title)
    title_font_size = min(28, max(18, int(28 * (80 / max(80, title_length)))))
    
    # Add title with auto-sizing text frame
    left = Inches(0.2)
    top = Inches(0)
    width = Inches(9)  # Full width
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Enable auto-sizing
    
    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,  # Use dynamic font size
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Set bullet points with proper spacing and formatting
    if isinstance(bullets, str):
        bullets = [bullets]
    
    bullets_left = Inches(0.5)
    bullets_top = Inches(1.5)
    bullets_width = Inches(9.0)
    bullets_height = Inches(3.5)
    
    bullets_box = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
    text_frame = bullets_box.text_frame
    text_frame.word_wrap = True
    
    for bullet in bullets:
        bullet_text = clean_text_content(bullet)
        
        # Calculate bullet font size based on length
        bullet_length = len(bullet_text)
        bullet_font_size = min(20, max(14, int(20 * (100 / max(100, bullet_length)))))
        
        p = text_frame.add_paragraph()
        p.text = "• " + bullet_text  # Add bullet character
        p.level = 0  # First level indentation
        apply_text_formatting(
            p, 
            font_size=bullet_font_size,  # Use dynamic font size
            color=template['colors']['text'],
            font_name=template['font']
        )
        
        # Add spacing between bullet points
        p.space_after = Pt(12)  # Add space after each bullet point
    
    return slide

def create_quote_slide(presentation, content, template):
    """Create a slide with a quote that adapts to content length"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Get content with no truncation
    quote = clean_text_content(content.get("quote", "Quote goes here"))
    author = clean_text_content(content.get("author", "Author"))
    
    # Calculate font sizes based on content length
    quote_length = len(quote)
    author_length = len(author)
    
    quote_font_size = min(20, max(14, int(20 * (200 / max(200, quote_length)))))
    author_font_size = min(12, max(10, int(12 * (40 / max(40, author_length)))))
    
    # Add quote text with proper centering
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(2.5)
    
    quote_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = quote_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = f'"{quote}"'
    apply_text_formatting(
        p, 
        font_size=quote_font_size,  # Use dynamic font size
        color=template['colors']['primary'],
        font_name=template['font'],
        italic=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Add author with proper right alignment
    left = Inches(4.0)
    top = Inches(4.0)
    width = Inches(5.0)
    height = Inches(0.8)
    
    author_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = author_box.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = f'— {author}'
    apply_text_formatting(
        p, 
        font_size=author_font_size,  # Use dynamic font size
        color=template['colors']['secondary'],
        font_name=template['font'],
        alignment=PP_ALIGN.RIGHT
    )
    
    return slide

def create_image_and_paragraph_slide(presentation, content, template):
    """Create a slide with an image placeholder on the left and paragraph text on right with adaptive sizing"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Set the slide background color to template background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = template['colors']['background']
    
    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Title'))
    paragraph = clean_text_content(content.get('paragraph', ''))
    imageDescription = clean_text_content(content.get('imageDescription', 'Image'))
    
    # Calculate dynamic font sizes based on content length
    title_length = len(title)
    paragraph_length = len(paragraph)
    
    title_font_size = max(20, min(32, int(32 * (80 / max(80, title_length)))))
    paragraph_font_size = max(9, min(14, int(14 * (300 / max(300, paragraph_length)))))
    
    # Add title with auto-sizing text frame
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,  # Use calculated size
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Add image placeholder with same dimensions
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(4)
    height = Inches(3)
    
    img_placeholder = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = template['colors']['secondary']
    
    # Add image description text
    if imageDescription:
        img_desc_width = Inches(3.6)  # Slightly smaller than the image width
        img_desc_height = Inches(1.5)  # Adjust height as needed
        
        # Center the text box within the image
        img_desc_left = left + (width - img_desc_width) / 2
        img_desc_top = top + (height - img_desc_height) / 2
        
        img_desc_box = slide.shapes.add_textbox(
            img_desc_left, 
            img_desc_top, 
            img_desc_width, 
            img_desc_height
        )
        text_frame = img_desc_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.add_paragraph()
        p.text = imageDescription
        p.alignment = PP_ALIGN.CENTER
        
        # Apply formatting
        for run in p.runs:
            run.font.size = Pt(12)  # Fixed size for the description preview
            run.font.name = template['font']
            run.font.color.rgb = template['colors']['background']
    
    # Add paragraph text on the right
    if paragraph:
        para_left = Inches(5)
        para_top = Inches(1.5)
        para_width = Inches(4.5)
        
        # Make height adaptive based on content length
        para_height = Inches(max(3, min(6, 3 + (paragraph_length // 300))))
        
        para_box = slide.shapes.add_textbox(
            para_left, para_top, para_width, para_height
        )
        text_frame = para_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        p = text_frame.add_paragraph()
        p.text = paragraph
        apply_text_formatting(
            p, 
            font_size=paragraph_font_size,  # Use calculated size
            color=template['colors']['text'],
            font_name=template['font'],
            alignment=PP_ALIGN.LEFT
        )
    
    return slide

def create_two_column_slide(presentation, content, template):
    """Create a slide with row-based layout that adapts to content length"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Get content with no truncation
    title = clean_text_content(content.get('title', 'Title'))
    column1_title = clean_text_content(content.get('column1Title', 'Column 1'))
    column1_content = clean_text_content(content.get('column1Content', 'Column 1 content'))
    column2_title = clean_text_content(content.get('column2Title', 'Column 2'))
    column2_content = clean_text_content(content.get('column2Content', 'Column 2 content'))
    
    # Calculate dynamic font sizes based on content length
    title_length = len(title)
    col1_title_length = len(column1_title)
    col2_title_length = len(column2_title)
    col1_content_length = len(column1_content)
    col2_content_length = len(column2_content)
    
    title_font_size = max(16, min(20, int(20 * (80 / max(80, title_length)))))
    col1_title_font_size = max(14, min(18, int(18 * (40 / max(40, col1_title_length)))))
    col2_title_font_size = max(14, min(18, int(18 * (40 / max(40, col2_title_length)))))
    col1_content_font_size = max(10, min(14, int(14 * (300 / max(300, col1_content_length)))))
    col2_content_font_size = max(10, min(14, int(14 * (300 / max(300, col2_content_length)))))
    
    # Add main title - full width
    left = Inches(0.5)
    top = Inches(0)
    width = Inches(9)
    height = Inches(1.5)  # Increased for two lines
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = title
    apply_text_formatting(
        p, 
        font_size=title_font_size,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # First section title - full width
    left = Inches(0.5)
    top = Inches(0.9)
    width = Inches(4.25)  # Half width with small gap
    height = Inches(0.75)
    
    section1_title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = section1_title_box.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = column1_title
    apply_text_formatting(
        p, 
        font_size=col1_title_font_size,
        color=template['colors']['secondary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # First section content - full width
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(4.25)  # Half width with small gap
    
    # Make height adaptive based on content length
    height = Inches(max(1.5, min(4, 1.5 + (col1_content_length // 200))))
    
    section1_content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = section1_content_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = column1_content
    apply_text_formatting(
        p, 
        font_size=col1_content_font_size,
        color=template['colors']['text'],
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT
    )
    
    # Second section title - full width
    left = Inches(5.25)  # Starting after first column
    top = Inches(0.9)
    width = Inches(4.25)  # Half width with small gap
    height = Inches(0.75)
    
    section2_title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = section2_title_box.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = column2_title
    apply_text_formatting(
        p, 
        font_size=col2_title_font_size,
        color=template['colors']['secondary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Second section content - full width
    left = Inches(5.25)  # Starting after first column
    top = Inches(1.3)
    width = Inches(4.25)  # Half width with small gap
    
    # Make height adaptive based on content length
    height = Inches(max(1.5, min(4, 1.5 + (col2_content_length // 200))))
    
    section2_content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = section2_content_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    p = text_frame.add_paragraph()
    p.text = column2_content
    apply_text_formatting(
        p, 
        font_size=col2_content_font_size,
        color=template['colors']['text'],
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT
    )
    
    return slide# Responsive PPT Generator Code Updates


import tkinter as tk
from tkinter import filedialog

def export_pptx_local():
    data = request.json
    slides = data.get('slides', [])
    template_id = data.get('template')
    topic = data.get('topic', 'presentation')

    try:
        # Create a new Tkinter root for each request
        root = tk.Tk()
        root.withdraw()  # Hide the main window
        
        # Bring the dialog to the front on Windows
        root.attributes('-topmost', True)
        
        filename = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")],
            initialfile=f"{topic.replace(' ', '_')}.pptx",
            title="Save Presentation As"
        )
        
        # Destroy the root to clean up Tkinter resources
        root.destroy()

        if not filename:
            return jsonify({'error': 'Export cancelled by user'}), 400

        # Generate the PPTX file
        create_presentation(filename, slides, template_id)

        return jsonify({
            'message': f'Presentation saved to {filename}'
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Export failed: {str(e)}'}), 500